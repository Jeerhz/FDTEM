#!/usr/bin/env python3
"""Build the experiment deck: training MT metrics on longer text."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

FIGS = Path("/tmp/claude-679803/-home-abensale-FDTEM/49dda269-ec34-44a9-b6f4-2f5e9153d255/scratchpad/figs")
OUT = Path("/home/abensale/FDTEM/docs/FDTEM_length_training_experiment.pptx")

INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
AQUA = RGBColor(0x1B, 0xAF, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    return s


def text(s, txt, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         space=6, line=1.15):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        p.line_spacing = line
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT
    return tb


def title(s, t, sub=None):
    text(s, t, .7, .45, 12, .9, size=32, bold=True)
    if sub:
        text(s, sub, .7, 1.22, 12, .5, size=16, color=INK2)


def rule(s, y=1.16, x=.7, w=11.9, color=None):
    from pptx.enum.shapes import MSO_SHAPE
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                            Inches(w), Pt(2))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color or RGBColor(0xE3, 0xE2, 0xDE)
    ln.line.fill.background()
    ln.shadow.inherit = False


def pic(s, name, y, height=None, width=None):
    p = FIGS / name
    if height:
        sh = s.shapes.add_picture(str(p), 0, Inches(y), height=Inches(height))
    else:
        sh = s.shapes.add_picture(str(p), 0, Inches(y), width=Inches(width))
    sh.left = int((W - sh.width) / 2)
    return sh


def tiles(s, items, y=2.4, h=1.55, gap=.28, color=BLUE):
    """Row of stat tiles: [(value, label), ...]"""
    from pptx.enum.shapes import MSO_SHAPE
    n = len(items)
    total = 11.9
    w = (total - gap * (n - 1)) / n
    for i, (val, lab) in enumerate(items):
        x = .7 + i * (w + gap)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                                 Inches(y), Inches(w), Inches(h))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF3)
        box.line.fill.background()
        box.shadow.inherit = False
        box.adjustments[0] = .06
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(.12)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        r.font.size = Pt(30); r.font.bold = True
        r.font.color.rgb = color; r.font.name = FONT
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = lab
        r2.font.size = Pt(13); r2.font.color.rgb = INK2; r2.font.name = FONT


# ── 1. title ─────────────────────────────────────────────────────────────────
s = slide()
text(s, "Does training on longer text\nmake MT metrics better?", .9, 2.1, 11.5, 2.2,
     size=46, bold=True, line=1.05)
text(s, "Finetuning COMET and COMET-QE on WMT paragraph-level data", .9, 4.45,
     11.5, .5, size=20, color=INK2)
text(s, "36 training runs   ·   262,080 training examples   ·   held-out WMT + MetaDocEval",
     .9, 5.05, 11.5, .5, size=16, color=BLUE)
text(s, "FDTEM  ·  August 2026", .9, 6.4, 11.5, .4, size=13, color=INK2)

# ── 2. the problem ───────────────────────────────────────────────────────────
s = slide()
title(s, "The mismatch", "COMET learns on sentences. It is applied to documents.")
rule(s, 1.75)
pic(s, "fig_gap.png", 2.35, width=10.5)
text(s, "Nobody knows what this costs.  Nor what fixes it.", .7, 5.6, 11.9, .6,
     size=22, color=INK)
text(s, "We measure both.", .7, 6.25, 11.9, .5, size=22, bold=True, color=ORANGE)

# ── 3. expected conclusions ──────────────────────────────────────────────────
s = slide()
title(s, "What we expect to find", "Stated before the results. The experiment can refute it.")
rule(s, 1.75)
pic(s, "fig_hypothesis.png", 2.0, width=9.8)
text(s,
     "1.  Long-only training helps on long text — and hurts on sentences.\n"
     "2.  A small share of sentences (~20%) restores the sentence-level skill.\n"
     "3.  Real documents teach more than concatenated sentences.\n"
     "4.  Better correlation ≠ understanding discourse errors.",
     .7, 6.05, 11.9, 1.2, size=15, color=INK, space=2)

# ── 4. source A ──────────────────────────────────────────────────────────────
s = slide()
title(s, "Source A — real long text", "WMT25 General MT: humans scored whole documents.")
rule(s, 1.75)
tiles(s, [("2,625", "documents\nwith reference"), ("47,924", "document × system\nhuman scores"),
          ("13", "language pairs"), ("82", "words per document\n(average)")],
      y=2.35, h=1.9, color=AQUA)
text(s, "Protocol:  ESA, one score 0–100 per document, plus error spans.\n"
        "Languages we never had before:  cs→de, cs→uk, en→ar, en→ja, en→is …\n"
        "After the validation split and the 512-token limit:  42,558 training rows.",
     .7, 4.75, 11.9, 1.6, size=17, color=INK, space=8)

# ── 5. source B ──────────────────────────────────────────────────────────────
s = slide()
title(s, "Source B — sentences, aggregated",
      "WMT22 MQM sentences glued into windows. The score is averaged.")
rule(s, 1.75)
pic(s, "fig_aggregation.png", 2.05, height=3.5)
text(s, "67,574 scored segments  ·  en-de, en-ru, zh-en  ·  15 MT systems\n"
        "→ 160,377 windows.  A window keeps the same score scale as a sentence.",
     .7, 5.85, 11.9, 1.1, size=17, color=INK, space=6)

# ── 6. the pool ──────────────────────────────────────────────────────────────
s = slide()
title(s, "The training pool", "Three kinds of examples. One common format: src · mt · ref · score.")
rule(s, 1.75)
pic(s, "fig_pool.png", 2.1, height=3.9)
text(s, "Splits are document-disjoint.  Scores are normalised per language pair.\n"
        "Rank-based evaluation, so the normalisation changes nothing downstream.",
     .7, 6.15, 11.9, 1.0, size=15, color=INK2, space=4)

# ── 7. mixtures ──────────────────────────────────────────────────────────────
s = slide()
title(s, "The experimental lever — data mixtures",
      "Same size every time. Only the composition changes.")
rule(s, 1.75)
pic(s, "fig_mixes.png", 2.05, width=11.6)
text(s, "9 mixtures × 24,000 examples.  Sampled from one shuffle, seed 42, md5 recorded.\n"
        "Two extra mixtures isolate the long data: real documents only vs aggregated only.",
     .7, 6.45, 11.9, 1.0, size=15, color=INK, space=4)

# ── 8. what we finetune ──────────────────────────────────────────────────────
s = slide()
title(s, "What we finetune", "Two public metrics. Two encoder regimes. Nine mixtures.")
rule(s, 1.75)
pic(s, "fig_grid.png", 2.1, height=3.3)
text(s, "COMET-DA  uses the reference.      CometKiwi QE  does not — source and translation only.\n"
        "Frozen = only the scoring head learns.  Trained = the whole encoder adapts.\n"
        "Everything else is identical: same schedule, same seed, early stopping on validation.",
     .7, 5.75, 11.9, 1.4, size=16, color=INK, space=6)

# ── 9. evaluation ────────────────────────────────────────────────────────────
s = slide()
title(s, "How we evaluate", "Two questions. Two test sets. Neither is used in training.")
rule(s, 1.75)
pic(s, "fig_eval.png", 2.05, height=3.4)
text(s, "Does the score still track human judgement?  →  WMT23 + WMT24 paragraphs, per length.\n"
        "Does the metric see discourse errors?  →  MetaDocEval: original vs perturbed document.",
     .7, 5.85, 11.9, 1.1, size=17, color=INK, space=8)

# ── 10. clean separation ─────────────────────────────────────────────────────
s = slide()
title(s, "No contamination", "Training years and evaluation years do not overlap.")
rule(s, 1.75)
tiles(s, [("2022 + 2025", "TRAIN\nsentences · documents"),
          ("2023 + 2024", "TEST\nparagraph correlation"),
          ("WMT24++", "TEST\nMetaDocEval discourse")],
      y=2.4, h=2.0, color=BLUE)
text(s, "MetaDocEval is built on WMT24 documents.  WMT24 is therefore excluded from training.\n"
        "The held-out sets also add unseen language pairs:  en-es and ja-zh.\n"
        "Uncertainty is measured by bootstrap over documents, not over examples.",
     .7, 4.9, 11.9, 1.6, size=17, color=INK, space=8)

# ── 11. status ───────────────────────────────────────────────────────────────
s = slide()
title(s, "Status", "Running on the cluster.")
rule(s, 1.75)
tiles(s, [("2", "runs finished"), ("8", "running now"),
          ("26", "queued"), ("0", "failures")], y=2.35, h=1.7, color=BLUE)
text(s, "First signal, preliminary:  the long-only COMET-DA run improves its validation\n"
        "correlation from 0.285 to 0.296.  Two runs out of 36 — not yet a result.",
     .7, 4.4, 11.9, 1.1, size=18, color=INK, space=8)
text(s, "Next:  full length-profile table per mixture, restoration point, frozen vs trained,\n"
        "DA vs QE, real documents vs aggregated sentences, then MetaDocEval accuracy.",
     .7, 5.7, 11.9, 1.1, size=16, color=INK2, space=8)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print("saved", OUT, f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
