#!/usr/bin/env python3
"""
make_results_pptx.py — build a PowerPoint deck of the FDTEM results.

Reads results/block_xsim/block_xsim.json + its plots and emits a self-contained
.pptx (16:9). Results-focused, with the block-xSIM++ recipe slides the study
hinges on. Run:  python scripts/make_results_pptx.py [--out docs/FDTEM_results.pptx]
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ── palette ───────────────────────────────────────────────────────────────────
INK   = RGBColor(0x0F, 0x1B, 0x2D)
MUTED = RGBColor(0x5A, 0x6B, 0x7B)
FAINT = RGBColor(0x94, 0xA2, 0xB0)
HAIR  = RGBColor(0xDB, 0xE2, 0xEA)
BG    = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF4, 0xF7, 0xFA)
TEAL  = RGBColor(0x0E, 0x7C, 0x86)   # Experiment A
TEALS = RGBColor(0xDA, 0xEC, 0xEE)
RUST  = RGBColor(0xC2, 0x57, 0x0C)   # Experiment B
RUSTS = RGBColor(0xF6, 0xE2, 0xD2)
GOOD  = RGBColor(0x15, 0x73, 0x47)
GOODS = RGBColor(0xDC, 0xEE, 0xE2)
GREY  = RGBColor(0x64, 0x74, 0x8B)
GREYS = RGBColor(0xE6, 0xEA, 0xEF)
CAUS  = RGBColor(0x7C, 0x3A, 0xED)
ENT   = RGBColor(0x02, 0x84, 0xC7)
NUM   = RGBColor(0xB4, 0x53, 0x09)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
FONT_H = "Calibri"
PLOTS = Path("results/block_xsim/plots")
BLOCK_JSON = Path("results/block_xsim/block_xsim.json")

# ── results (mean over languages), read from the block-xSIM++ output ───────────
_NAMES = {"comet:wmt22-comet-da": "COMET", "comet:4yeqp7cn-last": "Bio-COMET",
          "labse": "LaBSE", "e5:multilingual-e5-base": "E5",
          "xlmr:xlm-roberta-large": "XLM-R (raw)"}


def _load_results():
    j = json.loads(BLOCK_JSON.read_text())
    out = {}
    for spec, disp in _NAMES.items():
        if spec not in j["encoders"]:
            continue
        mk = j["encoders"][spec]["_mean_by_k"]
        out[disp] = {k: {m: mk[k].get(m) for m in
                         ("xsim_err", "xsimpp_err", "detection_rate")} for k in mk}
    return out


RESULTS = _load_results()

prs = Presentation()
prs.slide_width = I(13.333)
prs.slide_height = I(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ── helpers ───────────────────────────────────────────────────────────────────
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def _set(run, size, color, bold=False, font=FONT, italic=False):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = font


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line=1.06, wrap=True):
    """runs: list of paragraphs; each paragraph is a list of (txt, size, color, bold, italic?, font?)."""
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", 0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line
        for seg in para:
            txt, size, color = seg[0], seg[1], seg[2]
            bold = seg[3] if len(seg) > 3 else False
            italic = seg[4] if len(seg) > 4 else False
            fnt = seg[5] if len(seg) > 5 else FONT
            r = p.add_run(); r.text = txt; _set(r, size, color, bold, fnt, italic)
    return tb


def box(s, l, t, w, h, fill=None, line=None, line_w=1.0, rounded=True, dash=None):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
        if dash:
            d = shp.line._get_or_add_ln(); pd = d.makeelement(qn('a:prstDash'), {'val': dash})
            d.append(pd)
    if rounded:
        try: shp.adjustments[0] = 0.10
        except Exception: pass
    return shp


def eyebrow(s, txt, accent=RUST):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(0.6), I(0.5), I(0.09), I(0.32))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    bar.shadow.inherit = False
    text(s, I(0.82), I(0.5), I(11.8), I(0.34),
         [[(txt.upper(), 12, MUTED, True)]], anchor=MSO_ANCHOR.MIDDLE)


def title(s, txt, accent=INK, size=30):
    text(s, I(0.6), I(0.9), I(12.1), I(1.0), [[(txt, size, accent, True, False, FONT_H)]])


def tag(s, l, t, txt, fg, bg, w=I(1.7)):
    p = box(s, l, t, w, I(0.34), fill=bg, line=None)
    text(s, l, t, w, I(0.34), [[(txt, 11, fg, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return p


def footer(s, n):
    text(s, I(11.9), I(7.06), I(1.2), I(0.3),
         [[(f"{n:02d} / {TOTAL:02d}", 10, FAINT, False, False, "Consolas")]],
         align=PP_ALIGN.RIGHT)
    text(s, I(0.6), I(7.06), I(6), I(0.3),
         [[("FDTEM · COMET on long text", 10, FAINT)]])


def cells(s, l, t, labels, border, cellw=I(0.62), h=I(0.5), gap=I(0.06),
          fill=WHITE, txtcolor=None):
    """A row of small labelled cells = a 'block'. Returns right edge x."""
    txtcolor = txtcolor or MUTED
    x = l
    for lab in labels:
        c = box(s, x, t, cellw, h, fill=fill, line=border, line_w=1.25)
        text(s, x, t, cellw, h, [[(lab, 9, txtcolor, True, False, "Consolas")]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += cellw + gap
    return x


# ════════════════════════════════════════════════════════════════════════════
TOTAL = 14

# 1 — TITLE ---------------------------------------------------------------------
s = slide()
box(s, 0, 0, I(0.28), SH, fill=RUST, rounded=False)
text(s, I(0.9), I(1.2), I(11.6), I(0.4),
     [[("FDTEM · QUALITY ESTIMATION · LENGTH BIAS", 13, MUTED, True)]])
text(s, I(0.9), I(1.75), I(11.6), I(2.1),
     [[("Does COMET still work on ", 40, INK, True),
       ("longer text?", 40, RUST, True)],
      [("Results — block-level xSIM++, and the retraining plan", 20, MUTED)]],
     line=1.05)
tag(s, I(0.9), I(4.3), "EXPERIMENT A · retrain COMET", TEAL, TEALS, w=I(4.2))
text(s, I(0.9), I(4.75), I(5.4), I(1.4),
     [[("Better encoder + paragraph-level data. Does quality "
        "prediction hold up as text grows? ", 13, INK)],
      [("Training on the cluster — figures pending.", 12, MUTED, False, True)]], line=1.12)
tag(s, I(6.9), I(4.3), "EXPERIMENT B · block xSIM++", RUST, RUSTS, w=I(4.2))
text(s, I(6.9), I(4.75), I(5.4), I(1.4),
     [[("Bury one perturbed sentence in a k-sentence block: can the "
        "encoder still find the clean block? ", 13, INK)],
      [("Complete — 5 encoders, k = 2…5.", 12, RUST, True)]], line=1.12)
text(s, I(0.9), I(6.55), I(11.6), I(0.4),
     [[("Data & models in ~/scratch · logged to Weights & Biases · "
        "FLORES+ (openlanguagedata/flores_plus)", 11, FAINT)]])
footer(s, 1)

# 2 — FRAMING -------------------------------------------------------------------
s = slide()
eyebrow(s, "Motivation", RUST)
title(s, "An aligner pulls together; a quality metric must also push apart")
text(s, I(0.6), I(1.9), I(6.4), I(4.6),
     [[("Aligners (LaBSE, LASER, E5) are trained to place a sentence and "
        "its translation ", 15, INK), ("near", 15, INK, True),
       (" each other.", 15, INK)],
      [("", 6, INK)],
      [("COMET has a harder job: it must ", 15, INK),
       ("also separate", 15, RUST, True),
       (" a correct translation from one that is nearly identical but "
        "contains an error.", 15, INK)],
      [("", 8, INK)],
      [("The concern: ", 15, RUST, True),
       ("COMET only saw short segments in training. On a paragraph, is a "
        "single localised error simply diluted away?", 15, INK)]], line=1.16)
# right panel — two forces
box(s, I(7.4), I(1.9), I(5.3), I(4.5), fill=PANEL, line=HAIR)
tag(s, I(7.7), I(2.2), "PULL TOGETHER  ✓", GOOD, GOODS, w=I(2.6))
text(s, I(7.7), I(2.7), I(4.7), I(0.5),
     [[("src ↔ correct translation → close", 13, INK)]])
tag(s, I(7.7), I(3.5), "PUSH APART  ✗", RUST, RUSTS, w=I(2.6))
text(s, I(7.7), I(4.0), I(4.7), I(0.7),
     [[("src ↔ translation-with-one-error → must stay farther than the "
        "correct one", 13, INK)]], line=1.12)
box(s, I(7.7), I(5.0), I(4.7), I(0.02), fill=HAIR, rounded=False)
text(s, I(7.7), I(5.15), I(4.7), I(1.0),
     [[("Experiment B measures the push-apart force directly, as text "
        "gets longer.", 12.5, MUTED, False, True)]], line=1.14)
footer(s, 2)

# 3 — EXP B OVERVIEW ------------------------------------------------------------
s = slide()
eyebrow(s, "Experiment B · overview", RUST)
title(s, "xSIM++ on concatenated blocks")
text(s, I(0.6), I(1.85), I(12.1), I(0.9),
     [[("Standard xSIM++ perturbs one sentence. We move it to blocks of k "
        "sentences: perturb ", 15, INK),
       ("one", 15, RUST, True),
       (" sentence, keep the other k−1 intact, and ask whether the encoder "
        "can still retrieve the clean block.", 15, INK)]], line=1.16)
box(s, I(0.6), I(2.95), I(5.9), I(1.7), fill=PANEL, line=HAIR)
text(s, I(0.85), I(3.15), I(5.4), I(1.3),
     [[("THE QUESTION", 11, RUST, True)],
      [("As k grows, one error is a smaller fraction of the block. How fast "
        "does the encoder's sensitivity dilute?", 14, INK)]], line=1.14)
box(s, I(6.8), I(2.95), I(5.9), I(1.7), fill=PANEL, line=HAIR)
text(s, I(7.05), I(3.15), I(5.4), I(1.3),
     [[("THE METHOD IN ONE LINE", 11, RUST, True)],
      [("Cross-lingual retrieval: query = source block; find its true "
        "translation block in a pool seeded with hard negatives.", 14, INK)]], line=1.14)
text(s, I(0.6), I(5.0), I(12.1), I(1.3),
     [[("Encoder zoo (no training): ", 14, INK, True),
       ("COMET · Bio-COMET · raw XLM-R · LaBSE · E5.   ", 14, INK)],
      [("Data: FLORES+ dev+devtest, cased languages de · es · fr · ru "
        "(zh/ja/th excluded — entity edits need letter case).", 13.5, MUTED)]], line=1.3)
footer(s, 3)

# 4 — FLORES+ AGGREGATION -------------------------------------------------------
s = slide()
eyebrow(s, "Experiment B · aggregating FLORES+", RUST)
title(s, "From FLORES+ sentences to parallel blocks")
text(s, I(0.6), I(1.8), I(6.2), I(3.4),
     [[("FLORES+ rows are in document order with a url column marking "
        "articles. We take ", 14, INK),
       ("non-overlapping", 14, INK, True),
       (" windows of k consecutive sentences within one article "
        "(stride = k), and apply the ", 14, INK),
       ("same window to every language", 14, INK, True),
       (" — blocks stay parallel.", 14, INK)],
      [("", 6, INK)],
      [("• Pool dev + devtest → ", 13.5, INK),
       ("2,009 sentences", 13.5, RUST, True),
       (" / language, 562 articles (2–6 sentences).", 13.5, INK)],
      [("• Partial tail windows discarded — only full k-blocks kept.", 13.5, INK)],
      [("• Larger k ⇒ fewer blocks; k = 5 is the practical ceiling.", 13.5, INK)]],
     line=1.18)
# diagram: EN + DE parallel, k=3 window highlighted
text(s, I(7.0), I(1.85), I(5.6), I(0.3), [[("ONE ARTICLE · k = 3 · parallel", 10.5, MUTED, True)]])
text(s, I(7.0), I(2.3), I(0.4), I(0.5), [[("EN", 10, MUTED, True)]], anchor=MSO_ANCHOR.MIDDLE)
xr = cells(s, I(7.5), I(2.3), ["s1", "s2", "s3"], RUST)
cells(s, xr + I(0.12), I(2.3), ["s4", "s5", "s6"], HAIR, txtcolor=FAINT)
text(s, I(7.0), I(2.95), I(0.4), I(0.5), [[("DE", 10, MUTED, True)]], anchor=MSO_ANCHOR.MIDDLE)
xr = cells(s, I(7.5), I(2.95), ["t1", "t2", "t3"], RUST)
cells(s, xr + I(0.12), I(2.95), ["t4", "t5", "t6"], HAIR, txtcolor=FAINT)
text(s, I(7.0), I(3.6), I(5.6), I(0.3),
     [[("→ aligned source block & target block", 11, MUTED, False, True)]])
# counts table
tbl_data = [("block size", "blocks / language"), ("k = 2", "846"), ("k = 3", "468"),
            ("k = 4", "276"), ("k = 5", "141")]
ty = I(4.2)
box(s, I(7.5), ty, I(4.5), I(2.05), fill=PANEL, line=HAIR)
for i, (a, b) in enumerate(tbl_data):
    yy = ty + Emu(int(I(0.1)) + i * int(I(0.38)))
    bold = i == 0
    col = MUTED if i == 0 else INK
    text(s, I(7.7), yy, I(2.4), I(0.36), [[(a, 12, col, bold)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, I(9.9), yy, I(1.9), I(0.36), [[(b, 12, col, bold, False, "Consolas")]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 4)

# 5 — THE POOL (KEY) ------------------------------------------------------------
s = slide()
eyebrow(s, "Experiment B · the retrieval pool — the core of the recipe", RUST)
title(s, "What is the query, and what is it retrieved against?", size=27)
# query
text(s, I(0.6), I(1.7), I(3.4), I(0.3), [[("QUERY — one source block (EN)", 10.5, MUTED, True)]])
cells(s, I(0.6), I(2.05), ["s1", "s2", "s3"], RUST, cellw=I(0.7), h=I(0.55))
text(s, I(0.6), I(2.75), I(3.6), I(0.6),
     [[("embed → nearest candidate by cosine similarity", 11, MUTED, False, True)]], line=1.1)
# arrow
a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, I(3.9), I(2.15), I(0.7), I(0.34))
a.fill.solid(); a.fill.fore_color.rgb = FAINT; a.line.fill.background(); a.shadow.inherit = False
# pool
text(s, I(4.8), I(1.7), I(7.8), I(0.3), [[("CANDIDATE POOL (target language)", 10.5, MUTED, True)]])
py = I(2.05)
# clean
box(s, I(4.8), py, I(7.7), I(0.66), fill=GOODS, line=GOOD, line_w=1.5)
cells(s, I(4.95), py + I(0.07), ["t1", "t2", "t3"], GOOD, cellw=I(0.6), h=I(0.52), fill=WHITE)
text(s, I(7.2), py, I(5.1), I(0.66), [[("✓ the answer — clean block", 12, GOOD, True)]],
     anchor=MSO_ANCHOR.MIDDLE)
# hard neg 1
py2 = py + I(0.78)
box(s, I(4.8), py2, I(7.7), I(0.66), fill=RUSTS, line=RUST, line_w=1.5)
cells(s, I(4.95), py2 + I(0.07), ["t1", "t2*", "t3"], RUST, cellw=I(0.6), h=I(0.52), fill=WHITE)
text(s, I(7.2), py2, I(5.1), I(0.66),
     [[("hard negative — sentence 2 perturbed", 12, RUST, True)]], anchor=MSO_ANCHOR.MIDDLE)
# hard neg 2
py3 = py2 + I(0.78)
box(s, I(4.8), py3, I(7.7), I(0.66), fill=RUSTS, line=RUST, line_w=1.5)
cells(s, I(4.95), py3 + I(0.07), ["t1*", "t2", "t3"], RUST, cellw=I(0.6), h=I(0.52), fill=WHITE)
text(s, I(7.2), py3, I(5.1), I(0.66),
     [[("hard negative — different sentence / edit", 12, RUST, True)]], anchor=MSO_ANCHOR.MIDDLE)
# other true blocks
py4 = py3 + I(0.78)
box(s, I(4.8), py4, I(7.7), I(0.66), fill=None, line=GREY, line_w=1.25, dash="dash")
cells(s, I(4.95), py4 + I(0.07), ["u1", "u2", "u3"], GREY, cellw=I(0.6), h=I(0.52), fill=WHITE)
text(s, I(7.2), py4, I(5.1), I(0.66),
     [[("every OTHER true block (distractors)", 12, GREY, True)]], anchor=MSO_ANCHOR.MIDDLE)
# callout
box(s, I(0.6), I(5.75), I(12.1), I(1.15), fill=RUSTS, line=None)
box(s, I(0.6), I(5.75), I(0.09), I(1.15), fill=RUST, rounded=False)
text(s, I(0.85), I(5.85), I(11.7), I(1.0),
     [[("Crucial: ", 13.5, RUST, True),
       ("hard negatives are built per block, from that block's own k "
        "translations — one sentence perturbed, the other k−1 byte-identical. "
        "They are not perturbations of other blocks. The whole-corpus pool = "
        "all clean target blocks (mutual distractors) + every block's own "
        "hard negatives.", 13.5, INK)]], line=1.16)
footer(s, 5)

# 6 — PERTURBATIONS + WHAT/HOW MANY --------------------------------------------
s = slide()
eyebrow(s, "Experiment B · the perturbations", RUST)
title(s, "One sentence, three edit types, every block")
# three category cards
cats = [("CAUSALITY", CAUS, "antonym · negation · modal strengthening",
         "…ne s'amélioreraient pas → …s'amélioreraient"),
        ("ENTITY", ENT, "swap a named-entity span from a corpus bank",
         "M. Pittman → M. Phéniciens"),
        ("NUMBER", NUM, "change a digit / ordinal value",
         "88 → 98   ·   first → sixth")]
cw = I(3.95); cx = I(0.6)
for name, col, desc, ex in cats:
    box(s, cx, I(1.85), cw, I(1.95), fill=WHITE, line=HAIR)
    box(s, cx, I(1.85), cw, I(0.08), fill=col, rounded=False)
    tag(s, cx + I(0.2), I(2.05), name, col, PANEL, w=I(1.7))
    text(s, cx + I(0.2), I(2.55), cw - I(0.4), I(0.7), [[(desc, 12.5, INK)]], line=1.12)
    text(s, cx + I(0.2), I(3.25), cw - I(0.4), I(0.5),
         [[(ex, 10.5, MUTED, False, False, "Consolas")]], line=1.1)
    cx += cw + I(0.18)
# per-block sweep + pool sizes
text(s, I(0.6), I(4.1), I(6.0), I(2.4),
     [[("For each k-block, hard negatives sweep:", 14, INK, True)],
      [("• each position p = 1…k  (perturb p, keep the rest)", 13, INK)],
      [("• each category — causality, entity, number", 13, INK)],
      [("• a few variants each (default 2, deduplicated)", 13, INK)],
      [("", 5, INK)],
      [("The other k−1 sentences are byte-identical to the clean block — "
        "that is what makes the negative hard, and harder as k grows.", 12.5, MUTED, False, True)]],
     line=1.22)
# pool size table
box(s, I(7.0), I(4.1), I(5.7), I(2.35), fill=PANEL, line=HAIR)
text(s, I(7.2), I(4.2), I(5.3), I(0.3), [[("DE · pool sizes (variants/pos = 2)", 10.5, RUST, True)]])
rows = [("k", "true", "hard neg", "pool"), ("2", "846", "6,379", "7,225"),
        ("3", "468", "5,287", "5,755"), ("4", "276", "4,172", "4,448"),
        ("5", "141", "2,611", "2,752")]
colx = [I(7.3), I(8.5), I(9.9), I(11.4)]
for i, row in enumerate(rows):
    yy = I(4.6) + Emu(i * int(I(0.34)))
    bold = i == 0; col = MUTED if i == 0 else INK
    for cxx, val in zip(colx, row):
        text(s, cxx, yy, I(1.25), I(0.32),
             [[(val, 11.5, col, bold, False, "Consolas")]],
             align=PP_ALIGN.RIGHT if val != row[0] else PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE)
footer(s, 6)

# 7 — HOW WE SCORE --------------------------------------------------------------
s = slide()
eyebrow(s, "Experiment B · how we score it", RUST)
title(s, "Retrieval rule & the metrics we read")
text(s, I(0.6), I(1.85), I(12.1), I(0.9),
     [[("Embed the query block and every candidate; retrieval picks the "
        "candidate of ", 14.5, INK),
       ("maximum cosine similarity", 14.5, INK, True),
       (" (absolute-margin rule). A block is correct only if that argmax is "
        "its own clean target block.", 14.5, INK)]], line=1.16)
metrics = [("xsim error", "over the clean-blocks-only pool (classic)"),
           ("xsim++ error", "over the full pool incl. hard negatives"),
           ("detection rate", "P[ cos(q, clean) > cos(q, perturbed) ]"),
           ("error typology", "misaligned vs which category won"),
           ("by position", "detection vs where the error sits"),
           ("margin", "gold − best-negative similarity")]
y = I(3.0)
for i, (m, d) in enumerate(metrics):
    col = i % 2
    xx = I(0.6) + Emu(col * int(I(6.1)))
    yy = y + Emu((i // 2) * int(I(1.15)))
    box(s, xx, yy, I(5.9), I(1.0), fill=PANEL, line=HAIR)
    text(s, xx + I(0.25), yy + I(0.13), I(5.5), I(0.4), [[(m, 15, RUST, True)]])
    text(s, xx + I(0.25), yy + I(0.55), I(5.5), I(0.4),
         [[(d, 12.5, MUTED, False, False, "Consolas")]])
footer(s, 7)

# 8 — RESULT: HEADLINE ----------------------------------------------------------
s = slide()
eyebrow(s, "Experiment B · result · H100, 13 min · mean over de·es·fr·ru", RUST)
title(s, "COMET aligns blocks perfectly — yet can't feel a buried error", size=26)
s.shapes.add_picture(str(PLOTS / "xsim_vs_xsimpp.png"), I(0.5), I(1.75), width=I(6.15))
s.shapes.add_picture(str(PLOTS / "detection_vs_length.png"), I(6.75), I(1.75), width=I(6.15))
text(s, I(0.5), I(5.55), I(6.15), I(0.6),
     [[("Dashed = xsim (clean pool) ≈ 0 for COMET/LaBSE/E5. Solid = xsim++ "
        "(with hard negatives): COMET climbs 0.72 → 0.98.", 11, MUTED)]], line=1.1)
text(s, I(6.75), I(5.55), I(6.1), I(0.6),
     [[("P[clean closer than perturbed]. COMET & Bio-COMET fall below chance "
        "from k = 3; LaBSE/E5 hold ~0.85 at k = 5.", 11, MUTED)]], line=1.1)
box(s, I(0.6), I(6.25), I(12.1), I(0.72), fill=RUSTS, line=None)
box(s, I(0.6), I(6.25), I(0.09), I(0.72), fill=RUST, rounded=False)
text(s, I(0.85), I(6.32), I(11.7), I(0.6),
     [[("Finding: ", 12.5, RUST, True),
       ("COMET's encoder is fooled by a single perturbed sentence 72% of the "
        "time at k=2, 98% at k=5 — and from k=3 ranks the corrupted block "
        "closer to the source than the clean one. Aligners resist this.", 12.5, INK)]],
     line=1.12)
footer(s, 8)

# 9 — RESULT: WHICH & WHERE -----------------------------------------------------
s = slide()
eyebrow(s, "Experiment B · result · which edits, which positions", RUST)
title(s, "What fools each encoder, and where in the block")
s.shapes.add_picture(str(PLOTS / "error_by_category.png"), I(0.5), I(1.62), width=I(12.3))
text(s, I(0.6), I(3.98), I(12.1), I(0.55),
     [[("xsim++ error with the pool restricted to one category. COMET & XLM-R: "
        "entity ≥ causality ≫ number.  Aligners (LaBSE/E5): the reverse — "
        "entity/number are easy, causality is hardest.", 11.5, MUTED)]], line=1.14)
s.shapes.add_picture(str(PLOTS / "detection_by_position.png"), I(0.6), I(4.68), width=I(3.35))
text(s, I(4.25), I(5.0), I(8.4), I(1.8),
     [[("Detection vs the position of the perturbed sentence in the block.", 13, INK, True)],
      [("", 5, INK)],
      [("Two clean bands — aligners ~0.85–0.94 (top), COMET / XLM-R "
        "~0.33–0.50 (bottom) — with a slight dip for errors in the middle and "
        "later sentences of the block.", 12.5, MUTED)]], line=1.18)
footer(s, 9)

# 10 — RESULTS TABLE ------------------------------------------------------------
s = slide()
eyebrow(s, "Experiment B · result · the numbers", RUST)
title(s, "xsim++ error and detection rate, mean over languages")
rows = ["COMET", "Bio-COMET", "LaBSE", "E5", "XLM-R (raw)"]
ks = ["2", "3", "4", "5"]
# table: encoder | xsim++ (k2..k5) | detection (k2..k5)
tx, ty = I(0.6), I(1.9)
ncol = 1 + 4 + 4
colw = [I(2.5)] + [I(1.16)] * 8
th = I(0.55)
# header group labels
text(s, tx + colw[0], ty - I(0.02), Emu(sum(int(c) for c in colw[1:5])), I(0.32),
     [[("xsim++ error  (↑ = worse)", 12, RUST, True)]], align=PP_ALIGN.CENTER)
text(s, tx + Emu(int(colw[0]) + sum(int(c) for c in colw[1:5])), ty - I(0.02),
     Emu(sum(int(c) for c in colw[5:9])), I(0.32),
     [[("detection rate  (↓ = worse; .50 = chance)", 12, TEAL, True)]], align=PP_ALIGN.CENTER)
# header row
hy = ty + I(0.34)
heads = ["encoder"] + [f"k={k}" for k in ks] + [f"k={k}" for k in ks]
x = tx
for j, hd in enumerate(heads):
    text(s, x, hy, colw[j], I(0.4), [[(hd, 12, MUTED, True)]],
         align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += colw[j]
box(s, tx, hy + I(0.42), Emu(sum(int(c) for c in colw)), I(0.02), fill=HAIR, rounded=False)
# body
for i, enc in enumerate(rows):
    ry = hy + I(0.5) + Emu(i * int(I(0.62)))
    if i % 2 == 0:
        box(s, tx, ry, Emu(sum(int(c) for c in colw)), I(0.58), fill=PANEL, line=None)
    accent = RUST if enc in ("COMET", "Bio-COMET") else (TEAL if enc in ("LaBSE", "E5") else GREY)
    text(s, tx + I(0.1), ry, colw[0], I(0.58), [[(enc, 13, accent, True)]], anchor=MSO_ANCHOR.MIDDLE)
    x = tx + colw[0]
    for k in ks:
        v = RESULTS[enc][k]["xsimpp_err"]
        text(s, x, ry, colw[1], I(0.58), [[(f"{v:.2f}", 12.5, INK, False, False, "Consolas")]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += colw[1]
    for k in ks:
        v = RESULTS[enc][k]["detection_rate"]
        below = v < 0.5
        text(s, x, ry, colw[1], I(0.58),
             [[(f"{v:.2f}", 12.5, RUST if below else INK, below, False, "Consolas")]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += colw[1]
text(s, I(0.6), I(6.35), I(12.1), I(0.5),
     [[("Detection below .50 (bold rust) = the encoder ranks the corrupted "
        "block closer than the clean one, on average. COMET/Bio-COMET cross "
        "that line at k = 3.", 12, MUTED, False, True)]], line=1.14)
footer(s, 10)

# 11 — EXP A OVERVIEW -----------------------------------------------------------
s = slide()
eyebrow(s, "Experiment A · retraining", TEAL)
title(s, "Retrain COMET — isolate the encoder from the data", accent=INK)
box(s, I(0.6), I(1.9), I(5.9), I(2.0), fill=WHITE, line=HAIR)
box(s, I(0.6), I(1.9), I(5.9), I(0.08), fill=TEAL, rounded=False)
text(s, I(0.85), I(2.15), I(5.4), I(1.6),
     [[("ARM 1 · DATA EFFECT", 12, TEAL, True)],
      [("Continue wmt22-comet-da on paragraph-level data. Same encoder, "
        "same head — only the training text gets longer.", 13.5, INK)],
      [("Is short-segment data the limitation?", 12, MUTED, False, True)]], line=1.16)
box(s, I(6.8), I(1.9), I(5.9), I(2.0), fill=WHITE, line=HAIR)
box(s, I(6.8), I(1.9), I(5.9), I(0.08), fill=TEAL, rounded=False)
text(s, I(7.05), I(2.15), I(5.4), I(1.6),
     [[("ARM 2 · ENCODER EFFECT", 12, TEAL, True)],
      [("Retrain COMET's head on a better-aligned encoder "
        "(microsoft/infoxlm-large), from raw weights, on the same data.", 13.5, INK)],
      [("Is the ageing XLM-R backbone the limitation?", 12, MUTED, False, True)]], line=1.16)
box(s, I(0.6), I(4.1), I(12.1), I(0.95), fill=TEALS, line=None)
box(s, I(0.6), I(4.1), I(0.09), I(0.95), fill=TEAL, rounded=False)
text(s, I(0.85), I(4.2), I(11.7), I(0.8),
     [[("Controlled comparison. ", 13, TEAL, True),
       ("Arm 2 is run twice on identical data — infoXLM vs a plain "
        "xlm-roberta-large baseline — so any difference is attributable to the "
        "encoder, not the recipe.", 13, INK)]], line=1.14)
text(s, I(0.6), I(5.3), I(12.1), I(1.2),
     [[("Task throughout: ", 13.5, INK, True),
       ("reference-based scalar regression (RegressionMetric), MSE loss. "
        "Trained: regression head + layer-attention pool + encoder (low LR, "
        "30% frozen). Data: length-balanced Bio-MQM windows, k ∈ {1,2,3,4,6}, "
        "equal count per length — 15,410 train / 46,882 val.", 13.5, INK)]], line=1.22)
footer(s, 11)

# 12 — EXP A DATA ---------------------------------------------------------------
s = slide()
eyebrow(s, "Experiment A · data", TEAL)
title(s, "Length-balanced paragraph mix from Bio-MQM", accent=INK)
text(s, I(0.6), I(1.85), I(6.2), I(3.6),
     [[("Consecutive same-document Bio-MQM segments form real paragraphs. "
        "Slide a window of k segments, concatenate src/mt/ref, score by the ", 14, INK),
       ("mean per-segment MQM penalty", 14, INK, True),
       (". At k=1 this equals the sentence-level score — all lengths share "
        "one scale.", 14, INK)],
      [("", 6, INK)],
      [("• Balanced: every k contributes the same number of train windows.", 13, INK)],
      [("• Train windows overlap (augmentation); val is non-overlapping.", 13, INK)],
      [("• Windows over 480 tokens dropped (512 encoder limit).", 13, INK)]], line=1.2)
# window diagram
text(s, I(7.0), I(1.85), I(5.6), I(0.3), [[("ONE DOCUMENT → WINDOWS", 10.5, TEAL, True)]])
cells(s, I(7.0), I(2.25), ["s1", "s2", "s3", "s4", "s5", "s6"], HAIR, cellw=I(0.62), h=I(0.5), txtcolor=FAINT)
text(s, I(7.0), I(2.95), I(0.8), I(0.4), [[("k=2", 10, MUTED, True, False, "Consolas")]])
x = I(7.75)
for pair in [["s1", "s2"], ["s3", "s4"], ["s5", "s6"]]:
    x = cells(s, x, I(2.9), pair, TEAL, cellw=I(0.55), h=I(0.5)) + I(0.1)
text(s, I(7.0), I(3.55), I(0.8), I(0.4), [[("k=3", 10, MUTED, True, False, "Consolas")]])
x = I(7.75)
for tr in [["s1", "s2", "s3"], ["s4", "s5", "s6"]]:
    x = cells(s, x, I(3.5), tr, TEAL, cellw=I(0.55), h=I(0.5)) + I(0.1)
text(s, I(7.0), I(4.2), I(5.6), I(0.7),
     [[("Each window → one example: score = mean penalty. Length grows "
        "~40 → ~226 tokens (k=1 → k=6).", 11.5, MUTED, False, True)]], line=1.14)
# metric chips
chips = [("15,410", "train windows"), ("46,882", "val windows"), ("k∈{1,2,3,4,6}", "equal share")]
cx = I(0.6)
for v, l in chips:
    box(s, cx, I(5.55), I(3.9), I(1.0), fill=PANEL, line=HAIR)
    text(s, cx + I(0.25), I(5.68), I(3.5), I(0.5), [[(v, 20, INK, True)]])
    text(s, cx + I(0.25), I(6.15), I(3.5), I(0.35), [[(l, 12, MUTED)]])
    cx += I(4.05)
footer(s, 12)

# 13 — EXP A STATUS -------------------------------------------------------------
s = slide()
eyebrow(s, "Experiment A · status", TEAL)
title(s, "Retraining runs in flight", accent=INK)
text(s, I(0.6), I(1.9), I(6.2), I(3.4),
     [[("Three COMET runs training on RTX 8000 (≤ 3 days), logging to W&B "
        "project comet-retrain:", 14, INK)],
      [("", 6, INK)],
      [("• paragraph-continue — wmt22-comet-da on the length-balanced mix", 13.5, INK)],
      [("• infoXLM-large — encoder swap", 13.5, INK)],
      [("• XLM-R baseline — same recipe, controlled comparison", 13.5, INK)],
      [("", 6, INK)],
      [("On completion, eval_length_correlation.py produces the "
        "Kendall-τ-vs-k figure — one line per model. A model that generalises "
        "to long text degrades less as k grows.", 13, MUTED)]], line=1.2)
box(s, I(7.2), I(2.1), I(5.5), I(3.4), fill=PANEL, line=HAIR)
text(s, I(7.2), I(3.5), I(5.5), I(0.6),
     [[("Kendall τ vs k, per model", 14, FAINT, True)]], align=PP_ALIGN.CENTER)
text(s, I(7.2), I(4.0), I(5.5), I(0.5),
     [[("— pending training —", 12, FAINT, False, True)]], align=PP_ALIGN.CENTER)
footer(s, 13)

# 14 — RECAP --------------------------------------------------------------------
s = slide()
eyebrow(s, "Recap", RUST)
title(s, "Two experiments, one question")
box(s, I(0.6), I(1.85), I(5.9), I(3.5), fill=WHITE, line=HAIR)
box(s, I(0.6), I(1.85), I(5.9), I(0.08), fill=TEAL, rounded=False)
text(s, I(0.85), I(2.1), I(5.4), I(3.1),
     [[("A · RETRAINING", 13, TEAL, True)],
      [("data:  Bio-MQM windows k∈{1,2,3,4,6}", 13, INK)],
      [("task:  reference-based regression (MSE)", 13, INK)],
      [("trained:  head + pool + encoder (low LR)", 13, INK)],
      [("varies:  encoder (infoXLM) · or data", 13, INK)],
      [("eval:  Kendall τ vs k", 13, INK)],
      [("status:  training", 12.5, MUTED, False, True)]], line=1.34)
box(s, I(6.8), I(1.85), I(5.9), I(3.5), fill=WHITE, line=HAIR)
box(s, I(6.8), I(1.85), I(5.9), I(0.08), fill=RUST, rounded=False)
text(s, I(7.05), I(2.1), I(5.4), I(3.1),
     [[("B · BLOCK xSIM++", 13, RUST, True)],
      [("data:  FLORES+, non-overlapping k-blocks", 13, INK)],
      [("task:  cross-lingual block retrieval", 13, INK)],
      [("negatives:  1 sentence perturbed, k−1 intact", 13, INK)],
      [("varies:  block length k = 2…5 · encoder zoo", 13, INK)],
      [("eval:  xsim vs xsim++, detection, dilution", 13, INK)],
      [("result:  COMET blind to buried errors ✓", 12.5, RUST, True)]], line=1.34)
box(s, I(0.6), I(5.65), I(12.1), I(1.05), fill=INK, line=None)
text(s, I(0.9), I(5.78), I(11.5), I(0.85),
     [[("Does COMET's quality signal — and its encoder's error-sensitivity — "
        "survive the move from sentences to paragraphs?", 16, WHITE, True)]],
     anchor=MSO_ANCHOR.MIDDLE, line=1.16)
footer(s, 14)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default="docs/FDTEM_results.pptx")
    a = ap.parse_args()
    Path(a.out).parent.mkdir(parents=True, exist_ok=True)
    prs.save(a.out)
    print(f"wrote {a.out}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
