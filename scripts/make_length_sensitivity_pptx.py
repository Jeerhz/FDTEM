#!/usr/bin/env python3
"""
make_length_sensitivity_pptx.py — a concise, methodology-first deck on the
block-level xSIM++ length-sensitivity experiment.

How sensitive is a sentence encoder to a *single* translation error once that
error is buried inside a longer block of concatenated FLORES+ sentences? We
concatenate k parallel sentences into a block, plant one localised error in one
sentence, and measure how fast the encoder stops noticing it as k grows.

Reads results/block_xsim/block_xsim.json + its plots, emits a 16:9 .pptx.
Run:  python scripts/make_length_sensitivity_pptx.py
      [--out docs/FDTEM_length_sensitivity.pptx]
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── palette ───────────────────────────────────────────────────────────────────
INK   = RGBColor(0x0F, 0x1B, 0x2D)
MUTED = RGBColor(0x5A, 0x6B, 0x7B)
FAINT = RGBColor(0x94, 0xA2, 0xB0)
HAIR  = RGBColor(0xDB, 0xE2, 0xEA)
BG    = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF4, 0xF7, 0xFA)
RUST  = RGBColor(0xC2, 0x57, 0x0C)   # primary accent
RUSTS = RGBColor(0xF6, 0xE2, 0xD2)
TEAL  = RGBColor(0x0E, 0x7C, 0x86)
TEALS = RGBColor(0xDA, 0xEC, 0xEE)
GOOD  = RGBColor(0x15, 0x73, 0x47)
GOODS = RGBColor(0xDC, 0xEE, 0xE2)
GREY  = RGBColor(0x64, 0x74, 0x8B)
CAUS  = RGBColor(0x7C, 0x3A, 0xED)
ENT   = RGBColor(0x02, 0x84, 0xC7)
NUM   = RGBColor(0xB4, 0x53, 0x09)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
MONO = "Consolas"
PLOTS = Path("results/block_xsim/plots")
BLOCK_JSON = Path("results/block_xsim/block_xsim.json")

_NAMES = {"comet:wmt22-comet-da": "COMET", "comet:4yeqp7cn-last": "Bio-COMET",
          "labse": "LaBSE", "e5:multilingual-e5-base": "E5",
          "xlmr:xlm-roberta-large": "XLM-R (raw)"}


def _load_results():
    j = json.loads(BLOCK_JSON.read_text())
    out = {}
    for spec, disp in _NAMES.items():
        if spec not in j["encoders"]:
            continue
        mk = j["encoders"][spec]["_mean_by_k"]
        out[disp] = {k: {m: mk[k].get(m) for m in
                         ("xsim_err", "xsimpp_err", "detection_rate")} for k in mk}
    return out


RESULTS = _load_results()

DUEL_JSON = Path("results/block_duel/block_duel.json")
DPLOTS = Path("results/block_duel/plots")


def _load_duel():
    j = json.loads(DUEL_JSON.read_text())
    return {disp: j["encoders"][spec]["_mean_by_k"]
            for spec, disp in _NAMES.items() if spec in j["encoders"]}


DUEL = _load_duel()

prs = Presentation()
prs.slide_width = I(13.333)
prs.slide_height = I(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
TOTAL = 18


# ── helpers ───────────────────────────────────────────────────────────────────
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def _set(run, size, color, bold=False, font=FONT, italic=False):
    run.font.size = Pt(size); run.font.color.rgb = color
    run.font.bold = bold; run.font.italic = italic; run.font.name = font


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line=1.06, wrap=True):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for m in ("left", "right", "top", "bottom"):
        setattr(tf, f"margin_{m}", 0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line
        for seg in para:
            txt, size, color = seg[0], seg[1], seg[2]
            bold = seg[3] if len(seg) > 3 else False
            italic = seg[4] if len(seg) > 4 else False
            fnt = seg[5] if len(seg) > 5 else FONT
            r = p.add_run(); r.text = txt; _set(r, size, color, bold, fnt, italic)
    return tb


def box(s, l, t, w, h, fill=None, line=None, line_w=1.0, rounded=True, dash=None):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
        if dash:
            d = shp.line._get_or_add_ln(); pd = d.makeelement(qn('a:prstDash'), {'val': dash})
            d.append(pd)
    if rounded:
        try: shp.adjustments[0] = 0.09
        except Exception: pass
    return shp


def eyebrow(s, txt, accent=RUST):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(0.6), I(0.5), I(0.09), I(0.32))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    bar.shadow.inherit = False
    text(s, I(0.82), I(0.5), I(11.8), I(0.34),
         [[(txt.upper(), 12, MUTED, True)]], anchor=MSO_ANCHOR.MIDDLE)


def title(s, txt, accent=INK, size=29):
    text(s, I(0.6), I(0.9), I(12.1), I(1.0), [[(txt, size, accent, True)]])


def tag(s, l, t, txt, fg, bg, w=I(1.7), size=11):
    box(s, l, t, w, I(0.34), fill=bg, line=None)
    text(s, l, t, w, I(0.34), [[(txt, size, fg, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def footer(s, n):
    text(s, I(11.9), I(7.06), I(1.2), I(0.3),
         [[(f"{n:02d} / {TOTAL:02d}", 10, FAINT, False, False, MONO)]],
         align=PP_ALIGN.RIGHT)
    text(s, I(0.6), I(7.06), I(9), I(0.3),
         [[("Encoder length-sensitivity · block-level xSIM++ on FLORES+", 10, FAINT)]])


def cells(s, l, t, labels, border, cellw=I(0.62), h=I(0.5), gap=I(0.06),
          fill=WHITE, txtcolor=None):
    txtcolor = txtcolor or MUTED
    x = l
    for lab in labels:
        box(s, x, t, cellw, h, fill=fill, line=border, line_w=1.25)
        text(s, x, t, cellw, h, [[(lab, 9, txtcolor, True, False, MONO)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += cellw + gap
    return x


def bullets(s, l, t, w, h, items, size=13.5, line=1.22, gap=4):
    runs = []
    for it in items:
        if it == "":
            runs.append([("", 5, INK)])
        elif isinstance(it, list):
            runs.append(it)
        else:
            runs.append([("•  ", size, RUST, True), (it, size, INK)])
    return text(s, l, t, w, h, runs, line=line, space_after=gap)


# ════════════════════════════════════════════════════════════════════════════
# 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s = slide()
box(s, 0, 0, I(0.28), SH, fill=RUST, rounded=False)
text(s, I(0.9), I(1.15), I(11.6), I(0.4),
     [[("ENCODER EVALUATION · LENGTH SENSITIVITY", 13, MUTED, True)]])
text(s, I(0.9), I(1.7), I(11.8), I(2.2),
     [[("Do encoders still feel a single error ", 36, INK, True),
       ("in longer text?", 36, RUST, True)],
      [("Block-level xSIM++: one planted error, diluted across k concatenated "
        "FLORES+ sentences", 18, MUTED)]], line=1.08)
box(s, I(0.9), I(4.15), I(11.5), I(1.55), fill=PANEL, line=HAIR)
text(s, I(1.15), I(4.35), I(11.0), I(1.2),
     [[("The one-line idea.  ", 15, RUST, True),
       ("Concatenate k parallel sentences into a block, corrupt exactly one of "
        "them with a minimal semantic edit, and test whether an encoder can "
        "still retrieve the clean translation over the corrupted one. Sweep "
        "k = 2…5 so that ", 15, INK),
       ("length is the only thing that changes.", 15, INK, True)]], line=1.2)
text(s, I(0.9), I(6.15), I(11.6), I(0.5),
     [[("FLORES+ dev+devtest · de · es · fr · ru · pivot en   |   encoders: "
        "COMET · Bio-COMET · XLM-R · LaBSE · E5   |   no training, retrieval only",
        12, FAINT)]])
footer(s, 1)

# ════════════════════════════════════════════════════════════════════════════
# 2 — WHY
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Why this experiment")
title(s, "Aligning is easy; noticing one buried error is the hard part")
text(s, I(0.6), I(1.85), I(6.5), I(4.5),
     [[("A cross-lingual encoder must do two things:", 15, INK)],
      [("", 5, INK)],
      [("1.  pull a sentence and its translation ", 15, INK),
       ("together", 15, GOOD, True), (";", 15, INK)],
      [("2.  push a ", 15, INK), ("wrong", 15, RUST, True),
       (" translation — one that is nearly identical but carries an error — "
        "farther away than the correct one.", 15, INK)],
      [("", 8, INK)],
      [("Aligners (LaBSE, E5) are trained only for (1). A quality-aware "
        "encoder like COMET must also do (2).", 14.5, INK)],
      [("", 8, INK)],
      [("The worry: ", 15, RUST, True),
       ("COMET only ever saw short segments in training. In a long paragraph, "
        "is a single localised error simply averaged away?", 15, INK)]], line=1.2)
box(s, I(7.45), I(1.9), I(5.3), I(4.4), fill=PANEL, line=HAIR)
tag(s, I(7.75), I(2.2), "PULL TOGETHER  ✓", GOOD, GOODS, w=I(2.9))
text(s, I(7.75), I(2.72), I(4.7), I(0.5),
     [[("source ↔ correct translation → close", 13, INK, False, False, MONO)]])
tag(s, I(7.75), I(3.55), "PUSH APART  ✗ (the test)", RUST, RUSTS, w=I(3.4))
text(s, I(7.75), I(4.07), I(4.7), I(0.8),
     [[("source ↔ translation-with-one-error", 13, INK, False, False, MONO)],
      [("→ must stay farther than the clean one", 13, INK, False, False, MONO)]],
     line=1.15)
box(s, I(7.75), I(5.15), I(4.7), I(0.02), fill=HAIR, rounded=False)
text(s, I(7.75), I(5.3), I(4.7), I(0.9),
     [[("This deck measures the push-apart force, and how it decays as the "
        "block grows.", 12.5, MUTED, False, True)]], line=1.16)
footer(s, 2)

# ════════════════════════════════════════════════════════════════════════════
# 3 — RECIPE AT A GLANCE
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "The recipe at a glance")
title(s, "A cross-lingual retrieval task, made harder one length at a time")
steps = [
    ("1", "CONCATENATE", "k consecutive FLORES+ sentences → one block. Same "
     "window in every language, so source and target blocks stay parallel."),
    ("2", "CORRUPT", "Copy each target block, change exactly one sentence with a "
     "minimal semantic edit → a hard negative. Repeat per position & edit type."),
    ("3", "RETRIEVE", "Query = source block. Pool = every clean target block + "
     "all hard negatives. Correct only if the clean gold block is nearest."),
    ("4", "SWEEP k", "Run k = 2,3,4,5. The edit is unchanged; only the amount of "
     "surrounding context grows. Watch sensitivity dilute."),
]
cw = I(2.95); cx = I(0.6)
for num, head, desc in steps:
    box(s, cx, I(1.95), cw, I(3.2), fill=WHITE, line=HAIR)
    box(s, cx, I(1.95), cw, I(0.08), fill=RUST, rounded=False)
    text(s, cx + I(0.22), I(2.2), I(0.9), I(0.7),
         [[(num, 30, RUSTS, True)]], anchor=MSO_ANCHOR.TOP)
    text(s, cx + I(0.22), I(2.95), cw - I(0.44), I(0.4), [[(head, 13.5, RUST, True)]])
    text(s, cx + I(0.22), I(3.4), cw - I(0.44), I(1.7), [[(desc, 12.5, INK)]], line=1.18)
    cx += cw + I(0.14)
box(s, I(0.6), I(5.4), I(12.1), I(1.05), fill=RUSTS, line=None)
box(s, I(0.6), I(5.4), I(0.09), I(1.05), fill=RUST, rounded=False)
text(s, I(0.85), I(5.55), I(11.7), I(0.85),
     [[("Read the whole deck as: ", 13.5, RUST, True),
       ("hold the error fixed, vary the length, and see when the encoder stops "
        "reacting. Everything that follows exists to make length the single "
        "independent variable.", 13.5, INK)]], line=1.16)
footer(s, 3)

# ════════════════════════════════════════════════════════════════════════════
# 4 — CONCATENATION, PRECISELY
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Step 1 · concatenation — precisely")
title(s, "How sentences become parallel blocks")
bullets(s, I(0.6), I(1.85), I(6.4), I(4.6), [
    [("Non-overlapping windows. ", 14, INK, True),
     ("Within one FLORES+ article (rows share a url), take k consecutive "
      "sentences with stride = k. Windows never overlap, so no two blocks are "
      "near-duplicates.", 14, INK)],
    "",
    [("Same window, every language. ", 14, INK, True),
     ("The identical row range is applied to en, de, es, fr, ru — the source "
      "block and each target block are exact translations of one another.", 14, INK)],
    "",
    [("Join rule. ", 14, INK, True),
     ("Sentences are joined with a single space (with no separator for "
      "space-less scripts zh/ja/th). Nothing else is inserted.", 14, INK)],
    "",
    [("Only full blocks. ", 14, INK, True),
     ("A leftover tail shorter than k, or a window that would cross an article "
      "boundary, is discarded.", 14, INK)],
], size=13.5, line=1.16)
# diagram
text(s, I(7.25), I(1.85), I(5.6), I(0.3), [[("ONE ARTICLE · k = 3 · parallel", 10.5, MUTED, True)]])
text(s, I(7.25), I(2.28), I(0.4), I(0.5), [[("EN", 10, MUTED, True)]], anchor=MSO_ANCHOR.MIDDLE)
xr = cells(s, I(7.75), I(2.28), ["s1", "s2", "s3"], RUST)
cells(s, xr + I(0.12), I(2.28), ["s4", "s5", "s6"], HAIR, txtcolor=FAINT)
text(s, I(7.25), I(2.93), I(0.4), I(0.5), [[("DE", 10, MUTED, True)]], anchor=MSO_ANCHOR.MIDDLE)
xr = cells(s, I(7.75), I(2.93), ["t1", "t2", "t3"], RUST)
cells(s, xr + I(0.12), I(2.93), ["t4", "t5", "t6"], HAIR, txtcolor=FAINT)
text(s, I(7.25), I(3.58), I(5.6), I(0.3),
     [[("→ query = «s1 s2 s3»    gold = «t1 t2 t3»", 11, MUTED, False, True, MONO)]])
# counts
box(s, I(7.75), I(4.15), I(4.5), I(2.1), fill=PANEL, line=HAIR)
text(s, I(7.95), I(4.27), I(4.1), I(0.3), [[("blocks per language (de, dev+devtest)", 10.5, RUST, True)]])
tbl = [("k = 2", "846"), ("k = 3", "468"), ("k = 4", "276"), ("k = 5", "141")]
for i, (a, b) in enumerate(tbl):
    yy = I(4.62) + Emu(i * int(I(0.38)))
    text(s, I(7.95), yy, I(2.0), I(0.34), [[(a, 12.5, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, I(10.0), yy, I(2.0), I(0.34), [[(b, 12.5, INK, False, False, MONO)]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 4)

# ════════════════════════════════════════════════════════════════════════════
# 5 — THE ERRORS, PRECISELY
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Step 2 · the errors — precisely")
title(s, "Three edit types, each a minimal meaning change")
text(s, I(0.6), I(1.75), I(12.1), I(0.6),
     [[("The three perturbation families of xSIM++ (§2.2). Each edit changes "
        "one sentence's meaning with the smallest possible surface change, so "
        "the block stays fluent and near-identical.", 13.5, INK)]], line=1.16)
cats = [("CAUSALITY", CAUS,
         "antonym swap · add/remove a negation · strengthen a modal "
         "(may→will)",
         "«…ne s'amélioreraient pas» → «…s'amélioreraient»",
         "«augmenter de 12 %» → «diminuer de 12 %»"),
        ("ENTITY", ENT,
         "replace a detected named-entity span with another of the same width, "
         "sampled from a corpus-wide bank",
         "«chercheurs de Stanford» → «chercheurs de Pittman»", ""),
        ("NUMBER", NUM,
         "replace a digit with a random different value of the same length, or "
         "swap an ordinal",
         "«88 personnes» → «89 personnes»",
         "«première réunion» → «dixième réunion»")]
cw = I(3.95); cx = I(0.6)
for name, col, desc, ex1, ex2 in cats:
    box(s, cx, I(2.5), cw, I(2.55), fill=WHITE, line=HAIR)
    box(s, cx, I(2.5), cw, I(0.08), fill=col, rounded=False)
    tag(s, cx + I(0.2), I(2.7), name, col, PANEL, w=I(1.9))
    text(s, cx + I(0.2), I(3.2), cw - I(0.4), I(1.0), [[(desc, 12.5, INK)]], line=1.14)
    ey = I(4.25)
    text(s, cx + I(0.2), ey, cw - I(0.4), I(0.4),
         [[(ex1, 10, MUTED, False, False, MONO)]], line=1.1)
    if ex2:
        text(s, cx + I(0.2), ey + I(0.4), cw - I(0.4), I(0.4),
             [[(ex2, 10, MUTED, False, False, MONO)]], line=1.1)
    cx += cw + I(0.18)
box(s, I(0.6), I(5.35), I(12.1), I(1.15), fill=PANEL, line=HAIR)
text(s, I(0.85), I(5.48), I(11.7), I(1.0),
     [[("Which sentence, which side?  ", 13, RUST, True),
       ("Edits are applied to the pool side — the translations (direction "
        "en2xx) — to exactly one sentence at one position; the other k−1 "
        "sentences are left byte-identical. All edits are deterministic: the "
        "RNG is seeded from (seed, language, category, sentence, variant "
        "index), so the whole pool is reproducible.", 13, INK)]], line=1.16)
footer(s, 5)

# ════════════════════════════════════════════════════════════════════════════
# 6 — ONE ERROR PER BLOCK? + THE POOL
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Step 3 · the pool — how many errors per block?")
title(s, "One error per candidate — many single-error candidates per block")
# left: the pool illustration
text(s, I(0.6), I(1.8), I(6.3), I(0.3), [[("CANDIDATE POOL for one block (k = 3)", 10.5, MUTED, True)]])
py = I(2.15)
box(s, I(0.6), py, I(6.3), I(0.62), fill=GOODS, line=GOOD, line_w=1.5)
cells(s, I(0.75), py + I(0.06), ["t1", "t2", "t3"], GOOD, cellw=I(0.55), h=I(0.5), fill=WHITE)
text(s, I(2.9), py, I(3.9), I(0.62), [[("✓ gold — 0 errors", 12, GOOD, True)]], anchor=MSO_ANCHOR.MIDDLE)
rows = [(["t1*", "t2", "t3"], "pos 1 perturbed"),
        (["t1", "t2*", "t3"], "pos 2 perturbed"),
        (["t1", "t2", "t3*"], "pos 3 perturbed")]
for i, (lab, cap) in enumerate(rows):
    yy = py + Emu(int(I(0.72)) * (i + 1))
    box(s, I(0.6), yy, I(6.3), I(0.62), fill=RUSTS, line=RUST, line_w=1.25)
    cells(s, I(0.75), yy + I(0.06), lab, RUST, cellw=I(0.55), h=I(0.5), fill=WHITE)
    text(s, I(2.9), yy, I(3.9), I(0.62),
         [[(f"1 error — {cap}", 11.5, RUST, True)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, I(0.6), py + I(3.05), I(6.3), I(0.6),
     [[("× 3 edit types × up to 2 variants each → several single-error "
        "negatives per block.", 11, MUTED, False, True)]], line=1.12)
# right: the precise statement
bullets(s, I(7.2), I(1.9), I(5.5), I(4.4), [
    [("Each candidate carries exactly one error.", 13.5, INK, True)],
    [("A hard negative is a copy of the gold block with a single sentence "
      "edited — never two.", 13, INK)],
    "",
    [("A block contributes many such negatives.", 13.5, INK, True)],
    [("One per (position × edit type × variant), so longer blocks yield more "
      "negatives, not more errors per negative.", 13, INK)],
    "",
    [("The gold block has zero errors", 13, INK), (" and appears once.", 13, INK)],
    "",
    [("Whole pool = ", 13, INK, True),
     ("all clean target blocks (they distract each other) + every block's own "
      "single-error negatives.", 13, INK)],
], size=13, line=1.16)
footer(s, 6)

# ════════════════════════════════════════════════════════════════════════════
# 7 — LENGTH IS THE ONLY VARIABLE
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "The controls — why length is the only variable")
title(s, "Holding difficulty constant across block lengths")
ctrls = [
    ("PARALLEL BLOCKS",
     "The same sentence window is used in every language, so a source block and "
     "its gold target are exact translations — retrieval difficulty is not "
     "shifted by content, only by length."),
    ("ONE FIXED EDIT",
     "Every negative differs from its gold by exactly one localised edit; the "
     "other k−1 sentences are byte-identical. The error's type and size are the "
     "same at k = 2 and k = 5 — only its share of the block shrinks."),
    ("SAME EDIT BUDGET",
     "The same three categories and the same variants-per-position are applied "
     "at every k. What grows with k is the surrounding clean context, i.e. the "
     "dilution — nothing about the error itself."),
    ("PAIRWISE READOUT",
     "The headline metric — detection rate — compares the gold block against "
     "each of its own single-error negatives one-to-one. Because it is "
     "pairwise, a larger pool cannot confound it: it isolates dilution."),
]
cw = I(5.95); gap = I(0.2); cx0 = I(0.6)
for i, (h, d) in enumerate(ctrls):
    col = i % 2; row = i // 2
    xx = cx0 + Emu(col * (int(cw) + int(gap)))
    yy = I(1.95) + Emu(row * int(I(2.15)))
    box(s, xx, yy, cw, I(1.95), fill=WHITE, line=HAIR)
    box(s, xx, yy, I(0.09), I(1.95), fill=RUST, rounded=False)
    text(s, xx + I(0.28), yy + I(0.18), cw - I(0.5), I(0.4), [[(h, 14, RUST, True)]])
    text(s, xx + I(0.28), yy + I(0.62), cw - I(0.5), I(1.25), [[(d, 12.5, INK)]], line=1.16)
footer(s, 7)

# ════════════════════════════════════════════════════════════════════════════
# 8 — TWO GUARDS (edge cases)
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Edge cases — two guards that keep the pool clean")
title(s, "When an edit can't be made, and never a duplicate of the answer")
# guard 1
box(s, I(0.6), I(1.95), I(5.95), I(3.9), fill=WHITE, line=HAIR)
box(s, I(0.6), I(1.95), I(5.95), I(0.08), fill=TEAL, rounded=False)
tag(s, I(0.85), I(2.2), "GUARD 1 · edit not applicable", TEAL, TEALS, w=I(3.9), size=11.5)
bullets(s, I(0.85), I(2.75), I(5.5), I(3.0), [
    [("If a sentence has no number, no detectable entity, or no lexicon match, "
      "the perturbation function returns nothing.", 13, INK)],
    "",
    [("That (position, category) simply produces no negative. The error is "
      "never forced onto a sentence that cannot host it.", 13, INK)],
    "",
    [("So blocks do not all have the same negative count; per-category coverage "
      "is recorded in ", 13, INK),
     ("variant_stats", 12.5, MUTED, False, False, MONO),
     (" so any gap is explicit.", 13, INK)],
], size=13, line=1.18)
# guard 2
box(s, I(6.75), I(1.95), I(5.95), I(3.9), fill=WHITE, line=HAIR)
box(s, I(6.75), I(1.95), I(5.95), I(0.08), fill=GOOD, rounded=False)
tag(s, I(7.0), I(2.2), "GUARD 2 · never the gold in disguise", GOOD, GOODS, w=I(4.5), size=11.5)
bullets(s, I(7.0), I(2.75), I(5.5), I(3.0), [
    [("A candidate is kept only if the edited sentence actually differs from "
      "the original ", 13, INK),
     ("(cand ≠ source)", 12.5, MUTED, False, False, MONO),
     (", and duplicates are removed.", 13, INK)],
    "",
    [("A no-op edit therefore never enters the pool. The unperturbed block is "
      "never added as a negative — it would be identical to the gold and "
      "corrupt the labels.", 13, GOOD, True)],
    "",
    [("Generation over-samples (up to 6× draws) to reach the target variant "
      "count while discarding no-ops and repeats.", 13, INK)],
], size=13, line=1.18)
box(s, I(0.6), I(6.05), I(12.1), I(0.72), fill=PANEL, line=HAIR)
text(s, I(0.85), I(6.16), I(11.7), I(0.55),
     [[("Net effect: ", 12.5, RUST, True),
       ("every negative in the pool is a genuine, single-sentence corruption of "
        "its own block — no forced edits, no accidental copies of the correct "
        "answer.", 12.5, INK)]], line=1.14)
footer(s, 8)

# ════════════════════════════════════════════════════════════════════════════
# 9 — SCORING
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Step 3 · scoring")
title(s, "Retrieval rule and the metrics we read")
text(s, I(0.6), I(1.85), I(12.1), I(0.85),
     [[("Embed the query block and every candidate; retrieval picks the "
        "candidate of ", 14.5, INK),
       ("maximum cosine similarity", 14.5, INK, True),
       (" (absolute-margin rule). A block is scored correct only if that "
        "nearest candidate is its own clean gold block.", 14.5, INK)]], line=1.16)
metrics = [
    ("xsim error", "clean-blocks-only pool — pure alignment (should be ≈ 0)"),
    ("xsim++ error", "full pool incl. hard negatives — the aggregate task"),
    ("detection rate ★", "P[ cos(q, clean) > cos(q, perturbed) ] — length-controlled"),
    ("by position", "detection vs where in the block the error sits"),
    ("by category", "which edit type the encoder is blind to"),
    ("margin", "gold − best-negative cosine (how decisive)"),
]
y = I(3.0)
for i, (m, d) in enumerate(metrics):
    col = i % 2
    xx = I(0.6) + Emu(col * int(I(6.1)))
    yy = y + Emu((i // 2) * int(I(1.12)))
    star = "★" in m
    box(s, xx, yy, I(5.9), I(0.98), fill=(RUSTS if star else PANEL), line=HAIR)
    text(s, xx + I(0.25), yy + I(0.13), I(5.5), I(0.4), [[(m, 15, RUST, True)]])
    text(s, xx + I(0.25), yy + I(0.55), I(5.5), I(0.35),
         [[(d, 12, MUTED, False, False, MONO)]])
text(s, I(0.6), I(6.5), I(12.1), I(0.4),
     [[("★ detection rate is the strictly length-controlled number: a pairwise "
        "gold-vs-own-negative vote, immune to pool size. 0.50 = chance.",
        12, MUTED, False, True)]], line=1.12)
footer(s, 9)

# ════════════════════════════════════════════════════════════════════════════
# 10 — RESULT: HEADLINE
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Result · mean over de · es · fr · ru")
title(s, "COMET aligns blocks perfectly — yet can't feel a buried error", size=25)
s.shapes.add_picture(str(PLOTS / "xsim_vs_xsimpp.png"), I(0.5), I(1.7), width=I(6.15))
s.shapes.add_picture(str(PLOTS / "detection_vs_length.png"), I(6.75), I(1.7), width=I(6.15))
text(s, I(0.5), I(5.5), I(6.15), I(0.65),
     [[("Dashed xsim (clean pool) ≈ 0 for COMET/LaBSE/E5. Solid xsim++ (with "
        "hard negatives): COMET climbs 0.72 → 0.98 as k grows.", 11, MUTED)]], line=1.1)
text(s, I(6.75), I(5.5), I(6.1), I(0.65),
     [[("Detection = P[clean closer than perturbed]. COMET & Bio-COMET fall "
        "below chance from k = 3; aligners hold ~0.85 at k = 5.", 11, MUTED)]], line=1.1)
box(s, I(0.6), I(6.2), I(12.1), I(0.75), fill=RUSTS, line=None)
box(s, I(0.6), I(6.2), I(0.09), I(0.75), fill=RUST, rounded=False)
text(s, I(0.85), I(6.28), I(11.7), I(0.6),
     [[("Finding: ", 12.5, RUST, True),
       ("COMET's encoder aligns the block flawlessly but is fooled by one "
        "perturbed sentence 72 % of the time at k = 2 and 98 % at k = 5 — from "
        "k = 3 it ranks the corrupted block closer than the clean one. The pure "
        "aligners resist.", 12.5, INK)]], line=1.14)
footer(s, 10)

# ════════════════════════════════════════════════════════════════════════════
# 11 — RESULT: NUMBERS
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Result · the numbers")
title(s, "xsim++ error and detection rate vs block length")
enc_rows = ["COMET", "Bio-COMET", "LaBSE", "E5", "XLM-R (raw)"]
ks = ["2", "3", "4", "5"]
tx, ty = I(0.6), I(1.95)
colw = [I(2.5)] + [I(1.16)] * 8
text(s, tx + colw[0], ty, Emu(sum(int(c) for c in colw[1:5])), I(0.32),
     [[("xsim++ error   (↑ = worse)", 12, RUST, True)]], align=PP_ALIGN.CENTER)
text(s, tx + Emu(int(colw[0]) + sum(int(c) for c in colw[1:5])), ty,
     Emu(sum(int(c) for c in colw[5:9])), I(0.32),
     [[("detection rate   (↓ = worse; .50 = chance)", 12, TEAL, True)]], align=PP_ALIGN.CENTER)
hy = ty + I(0.36)
heads = ["encoder"] + [f"k={k}" for k in ks] + [f"k={k}" for k in ks]
x = tx
for j, hd in enumerate(heads):
    text(s, x, hy, colw[j], I(0.4), [[(hd, 12, MUTED, True)]],
         align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += colw[j]
box(s, tx, hy + I(0.42), Emu(sum(int(c) for c in colw)), I(0.02), fill=HAIR, rounded=False)
for i, enc in enumerate(enc_rows):
    ry = hy + I(0.5) + Emu(i * int(I(0.64)))
    if i % 2 == 0:
        box(s, tx, ry, Emu(sum(int(c) for c in colw)), I(0.6), fill=PANEL, line=None)
    accent = RUST if enc in ("COMET", "Bio-COMET") else (TEAL if enc in ("LaBSE", "E5") else GREY)
    text(s, tx + I(0.1), ry, colw[0], I(0.6), [[(enc, 13, accent, True)]], anchor=MSO_ANCHOR.MIDDLE)
    x = tx + colw[0]
    for k in ks:
        v = RESULTS[enc][k]["xsimpp_err"]
        text(s, x, ry, colw[1], I(0.6), [[(f"{v:.2f}", 12.5, INK, False, False, MONO)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += colw[1]
    for k in ks:
        v = RESULTS[enc][k]["detection_rate"]
        below = v is not None and v < 0.5
        txt = f"{v:.2f}" if v is not None else "—"
        text(s, x, ry, colw[1], I(0.6),
             [[(txt, 12.5, RUST if below else INK, below, False, MONO)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += colw[1]
text(s, I(0.6), I(6.35), I(12.1), I(0.55),
     [[("Bold rust = detection below .50: the encoder ranks the corrupted block "
        "closer than the clean one, on average. COMET and Bio-COMET cross that "
        "line at k = 3; LaBSE and E5 never do.", 12, MUTED, False, True)]], line=1.14)
footer(s, 11)

# ════════════════════════════════════════════════════════════════════════════
# 12 — RECAP
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Recap")
title(s, "One planted error, four lengths, one conclusion")
box(s, I(0.6), I(1.9), I(12.1), I(3.35), fill=WHITE, line=HAIR)
recap = [
    ("concatenate", "non-overlapping k-sentence windows, same window in every "
     "language → parallel source & target blocks (k = 2…5)."),
    ("corrupt", "one sentence per candidate, edited by one of three minimal "
     "semantic edits — causality, entity, number; the other k−1 untouched."),
    ("control", "the edit is fixed across k; only the surrounding context grows. "
     "Detection rate is a pairwise, pool-size-immune readout of dilution."),
    ("safeguards", "no edit is forced on a sentence that can't host it; a no-op "
     "is never added — the gold is never duplicated as a negative."),
    ("result", "COMET aligns blocks perfectly (xsim ≈ 0) but its error "
     "sensitivity dilutes fast: below chance from k = 3, 0.98 xsim++ at k = 5. "
     "Pure aligners (LaBSE, E5) resist far better."),
]
ry = I(2.12)
for rlab, rdesc in recap:
    tag(s, I(0.85), ry, rlab.upper(), RUST, RUSTS, w=I(1.75), size=10.5)
    text(s, I(2.8), ry - I(0.02), I(9.6), I(0.58), [[(rdesc, 12.5, INK)]], line=1.14)
    ry += I(0.62)
box(s, I(0.6), I(5.42), I(12.1), I(1.3), fill=RUSTS, line=None)
box(s, I(0.6), I(5.42), I(0.09), I(1.3), fill=RUST, rounded=False)
text(s, I(0.85), I(5.54), I(11.7), I(1.1),
     [[("Part II — fixed-pool duels (next slides). ", 12.5, RUST, True),
       ("Pin the pool to the gold + D single-error negatives of its own block "
        "at every k, averaged over all D-subsets — constant difficulty, no "
        "pool-composition confound. ", 12.5, INK),
       ("Coverage caveat: ", 12.5, RUST, True),
       ("D = 6 (the full k = 2 budget, 2 positions × 3 error types) requires "
        "both sentences of a k = 2 block to host all three edit types — only "
        "78 / 846 de blocks (≈ 9 %) qualify — so D = 5 and D = 4 are scored "
        "alongside from the same embeddings, trading pool size for coverage.",
        12.5, INK)]], line=1.16)
footer(s, 12)



# ════════════════════════════════════════════════════════════════════════════
# 13 — PART II · WHY D=6 COVERAGE COLLAPSES AT k=2
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Part II · fixed-pool duels — why 6-negative coverage is thin")
title(s, "A negative must be possible — six at once rarely are", size=27)
bullets(s, I(0.6), I(1.8), I(6.3), I(4.9), [
    [("A (position, category) pair yields a negative only if the edit can "
      "fire on that sentence AND change the string:", 13, INK)],
    "",
    [("NUMBER — ", 13, NUM, True),
     ("the sentence must contain a digit or an ordinal word; the redrawn "
      "value must differ.", 13, INK)],
    [("ENTITY — ", 13, ENT, True),
     ("needs a capitalised span (≥ 2 chars) that never appears lowercased in "
      "the corpus, plus a same-width replacement in the bank.", 13, INK)],
    [("CAUSALITY — ", 13, CAUS, True),
     ("needs an antonym-lexicon hit, a boostable modal (may→will), or a "
      "negation that can be added / removed.", 13, INK)],
    "",
    [("D = 6 at k = 2 demands all three categories on BOTH sentences — a "
      "conjunction of six events. FLORES sentences often lack digits, and "
      "es/fr/ru entity spans are rare, so the joint probability collapses.",
      13, INK, True)],
    "",
    [("German is the exception: every noun is capitalised, so the entity "
      "edit almost always applies — 78/846 blocks vs 8–9 elsewhere.", 13, INK)],
], size=13, line=1.16)
s.shapes.add_picture(str(DPLOTS / "duel_coverage.png"), I(7.15), I(1.85), width=I(5.6))
box(s, I(7.15), I(5.55), I(5.6), I(1.15), fill=PANEL, line=HAIR)
text(s, I(7.4), I(5.65), I(5.1), I(0.3),
     [[("blocks with ≥ D negatives at k = 2 (of 846)", 10.5, RUST, True)]])
rowsd = [("de", "78/293/843"), ("es", "9/62/220"),
         ("fr", "8/56/155"), ("ru", "8/46/158")]
for i, (lg, v3) in enumerate(rowsd):
    xx = I(7.4) + Emu(i * int(I(1.32)))
    text(s, xx, I(6.0), I(1.3), I(0.6),
         [[(lg + " ", 11.5, INK, True), (v3, 10, MUTED, False, False, MONO)]],
         line=1.1)
text(s, I(7.4), I(6.42), I(5.1), I(0.25),
     [[("per language: D = 6 / 5 / 4", 9.5, FAINT, False, True)]])
footer(s, 13)

# ════════════════════════════════════════════════════════════════════════════
# 14 — PART II · PROTOCOL & COUNTS
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Part II · the fixed-pool protocol")
title(s, "Same number of candidates at every k — the counts, exactly")
bullets(s, I(0.6), I(1.9), I(5.7), I(4.6), [
    [("One variant per (position, category)", 13.5, INK, True),
     (" → a k-sentence block owns at most 3k single-error negatives.",
      13.5, INK)],
    "",
    [("A duel is the gold + D of them", 13.5, INK, True),
     (" (D = 6, 5, 4). Score = mean retrieval accuracy over ALL C(m, D) "
      "subsets of the block's m valid negatives.", 13.5, INK)],
    "",
    [("Closed form, no enumeration: ", 13.5, INK, True),
     ("the gold wins a D-subset iff it beats every member, so", 13.5, INK)],
    [("P[win] = C(w, D) / C(m, D)", 12.5, RUST, True, False, MONO)],
    [("w = #negatives beaten — verified against brute-force enumeration.",
      12, MUTED, False, True)],
    "",
    [("Blocks with m < D are skipped", 13.5, INK, True),
     (" — every scored duel has exactly D + 1 candidates; a tie counts "
      "against the gold.", 13.5, INK)],
], size=13.5, line=1.18)
tx, ty = I(6.6), I(2.0)
colw = [I(0.7), I(1.2), I(1.5), I(1.0), I(1.0), I(1.0)]
heads = ["k", "blocks", "max m = 3k", "C(3k,6)", "C(3k,5)", "C(3k,4)"]
x = tx
for j, hd in enumerate(heads):
    text(s, x, ty, colw[j], I(0.4), [[(hd, 11.5, MUTED, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += colw[j]
box(s, tx, ty + I(0.44), Emu(sum(int(c) for c in colw)), I(0.02), fill=HAIR, rounded=False)
crow = [("2", "846", "6", "1", "6", "15"),
        ("3", "468", "9", "84", "126", "126"),
        ("4", "276", "12", "924", "792", "495"),
        ("5", "141", "15", "5005", "3003", "1365")]
for i, vals in enumerate(crow):
    ry = ty + I(0.52) + Emu(i * int(I(0.6)))
    if i % 2 == 0:
        box(s, tx, ry, Emu(sum(int(c) for c in colw)), I(0.56), fill=PANEL, line=None)
    x = tx
    for j, v in enumerate(vals):
        text(s, x, ry, colw[j], I(0.56),
             [[(v, 12.5, INK, j == 0, False, MONO)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += colw[j]
text(s, tx, ty + I(3.15), I(6.3), I(1.5),
     [[("blocks = per language (dev + devtest), identical across languages. "
        "In practice a block has m ≤ 3k valid negatives (conditions on the "
        "previous slide); the duel averages over the C(m, D) subsets that "
        "exist and skips blocks with m < D.", 11.5, MUTED)]], line=1.18)
footer(s, 14)

# ════════════════════════════════════════════════════════════════════════════
# 15 — PART II · RESULTS HEADLINE
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Part II · results — constant pool size, same conclusion")
title(s, "Dilution survives the pool-size control", size=27)
s.shapes.add_picture(str(DPLOTS / "duel_err_vs_length.png"), I(0.75), I(1.75), width=I(11.8))
text(s, I(0.75), I(5.62), I(11.8), I(0.5),
     [[("One panel per duel size; each has its own chance line D/(D+1): "
        "6/7 ≈ .857, 5/6 ≈ .833, 4/5 = .800. Encoder ranking and the upward "
        "dilution trend are identical in all three.", 11, MUTED)]], line=1.12)
box(s, I(0.6), I(6.2), I(12.1), I(0.75), fill=RUSTS, line=None)
box(s, I(0.6), I(6.2), I(0.09), I(0.75), fill=RUST, rounded=False)
text(s, I(0.85), I(6.28), I(11.7), I(0.6),
     [[("Finding: ", 12.5, RUST, True),
       ("COMET / Bio-COMET / raw XLM-R sit at or above chance from k = 2–3 "
        "even with the pool pinned to D + 1 candidates; LaBSE and E5 stay far "
        "below chance at every k. D = 4 is the recommended headline: chance "
        "drops only .857 → .800 while the k = 2 sample grows 103 → 1376 "
        "blocks (13×).", 12.5, INK)]], line=1.14)
footer(s, 15)

# ════════════════════════════════════════════════════════════════════════════
# 16 — PART II · RESULTS NUMBERS
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Part II · results — the numbers")
title(s, "Duel error vs block length (mean over de · es · fr · ru)")
enc_rows2 = ["COMET", "Bio-COMET", "LaBSE", "E5", "XLM-R (raw)"]
ks2 = ["2", "3", "4", "5"]
tx, ty = I(0.6), I(1.95)
colw = [I(2.5)] + [I(1.16)] * 8
text(s, tx + colw[0], ty, Emu(sum(int(c) for c in colw[1:5])), I(0.32),
     [[("D = 6   (chance .857)", 12, RUST, True)]], align=PP_ALIGN.CENTER)
text(s, tx + Emu(int(colw[0]) + sum(int(c) for c in colw[1:5])), ty,
     Emu(sum(int(c) for c in colw[5:9])), I(0.32),
     [[("D = 4   (chance .800)", 12, TEAL, True)]], align=PP_ALIGN.CENTER)
hy = ty + I(0.36)
heads = ["encoder"] + [f"k={k}" for k in ks2] + [f"k={k}" for k in ks2]
x = tx
for j, hd in enumerate(heads):
    text(s, x, hy, colw[j], I(0.4), [[(hd, 12, MUTED, True)]],
         align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += colw[j]
box(s, tx, hy + I(0.42), Emu(sum(int(c) for c in colw)), I(0.02), fill=HAIR, rounded=False)
for i, enc in enumerate(enc_rows2):
    ry = hy + I(0.5) + Emu(i * int(I(0.64)))
    if i % 2 == 0:
        box(s, tx, ry, Emu(sum(int(c) for c in colw)), I(0.6), fill=PANEL, line=None)
    accent = RUST if enc in ("COMET", "Bio-COMET") else (TEAL if enc in ("LaBSE", "E5") else GREY)
    text(s, tx + I(0.1), ry, colw[0], I(0.6), [[(enc, 13, accent, True)]], anchor=MSO_ANCHOR.MIDDLE)
    x = tx + colw[0]
    for D, chance in (("6", 6 / 7), ("4", 4 / 5)):
        for k in ks2:
            v = DUEL[enc][k]["by_duel_size"][D]["duel_err"]
            bad = v is not None and v >= chance
            vtxt = f"{v:.2f}" if v is not None else "—"
            text(s, x, ry, colw[1], I(0.6),
                 [[(vtxt, 12.5, RUST if bad else INK, bad, False, MONO)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x += colw[1]
text(s, I(0.6), I(6.15), I(12.1), I(0.8),
     [[("Bold rust = at or above the duel's chance line — no better than "
        "picking one of the D + 1 candidates at random. Coverage at k = 2 "
        "(blocks scored, of 3384): D = 6 → 103, D = 5 → 457, D = 4 → 1376. "
        "D = 5 sits between the two columns shown (full numbers in "
        "results/block_duel/block_duel.json).", 12, MUTED, False, True)]],
     line=1.14)
footer(s, 16)

# ════════════════════════════════════════════════════════════════════════════
# 17 — PART II · CATEGORY SENSITIVITY
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Part II · sensitivity by error category")
title(s, "Causality errors are the hardest — for everyone", size=27)
s.shapes.add_picture(str(DPLOTS / "beat_by_category.png"), I(0.55), I(1.8), width=I(12.2))
text(s, I(0.55), I(4.25), I(12.2), I(0.4),
     [[("P[gold beats the negative] per category (bars per k = 2,3,4,5; "
        "dotted line = chance .50). One panel per encoder.", 11, MUTED)]],
     line=1.1)
box(s, I(0.6), I(4.8), I(12.1), I(1.85), fill=WHITE, line=HAIR)
bullets(s, I(0.9), I(5.0), I(11.6), I(1.6), [
    [("Aligners: ", 13, TEAL, True),
     ("number and entity swaps stay easy at k = 5 (LaBSE .96 / .89, E5 .93 / "
      ".94) but causality falls to .68–.69 — antonyms and negation move the "
      "embedding least.", 13, INK)],
    [("COMET / Bio-COMET: ", 13, RUST, True),
     ("no category survives — all three are near .50 already at k = 2 "
      "(causality .56, entity .62, number .53) and sink to ~.33–.37 by "
      "k = 5. The blindness is category-independent.", 13, INK)],
], size=13, line=1.2)
footer(s, 17)

# ════════════════════════════════════════════════════════════════════════════
# 18 — PART II · POSITION SENSITIVITY
# ════════════════════════════════════════════════════════════════════════════
s = slide()
eyebrow(s, "Part II · sensitivity by error position")
title(s, "Where the error sits barely matters — depth does", size=27)
s.shapes.add_picture(str(DPLOTS / "beat_by_position.png"), I(0.6), I(1.85), width=I(6.6))
bullets(s, I(7.55), I(2.0), I(5.2), I(4.4), [
    [("Aligners show a mild late-position penalty. ", 13.5, TEAL, True),
     ("LaBSE at k = 5 detects .86 of first-sentence errors but .76 of "
      "last-sentence ones; E5 drifts .87 → .82. Errors buried deep in the "
      "block are slightly harder.", 13.5, INK)],
    "",
    [("COMET is uniformly blind. ", 13.5, RUST, True),
     ("Detection is flat across positions (~.45–.48 at k = 3, ~.30–.36 at "
      "k = 5): the failure is not a positional artefact — the whole block is "
      "averaged into one vector that one sentence cannot move.", 13.5, INK)],
    "",
    [("Takeaway: ", 13.5, INK, True),
     ("length itself, not error placement, drives the dilution.", 13.5, INK)],
], size=13.5, line=1.2)
footer(s, 18)


def main() -> None:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--out", default="docs/FDTEM_length_sensitivity.pptx")
    args = ap.parse_args()
    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"deck → {out}  ({TOTAL} slides)")


if __name__ == "__main__":
    main()
